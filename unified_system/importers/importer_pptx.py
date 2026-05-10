"""
importer_pptx.py
Nhập file PPTX vào bảng:
  bai_giang, trang, khoi_noi_dung, tai_nguyen_media,
  doan_van_trang, tu_khoa_trang, chu_de_trang, chu_de_bai_giang
"""
import os, sys, io, json, hashlib, re


from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv
import psycopg2
from psycopg2.extras import Json

load_dotenv(Path(__file__).parent.parent / ".env")

# ── Thư viện PPTX ────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Thiếu: pip install python-pptx")
    sys.exit(1)

# ── Kết nối DB ───────────────────────────────────────────────────
DB = dict(
    host=os.getenv("DB_HOST", "localhost"),
    port=os.getenv("DB_PORT", "5432"),
    dbname=os.getenv("DB_NAME", "kho_bai_giang"),
    user=os.getenv("DB_USER", "postgres"),
    password=os.getenv("DB_PASSWORD", "123456"),
)

# ── Nguồn gốc theo thư mục cha ───────────────────────────────────
def detect_nguon_goc(filepath: Path) -> str:
    parts = [p.lower() for p in filepath.parts]
    for p in parts:
        if "01_thu_thap_web" in p or "thu_thap" in p or "web" in p:
            return "thu_thap_web"
        if "02_ai_tao_sinh" in p or "ai_tao" in p or "gemini" in p or "canva" in p:
            return "ai_tao_sinh"
    return "thu_cong"

# ── Hash file ────────────────────────────────────────────────────
def file_hash(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()

# ── Trích xuất text từ shape ─────────────────────────────────────
def extract_shape_text(shape):
    """Trả về (full_text, list[dict paragraph])"""
    if not shape.has_text_frame:
        return "", []
    paragraphs = []
    all_lines = []
    for p_idx, para in enumerate(shape.text_frame.paragraphs):
        line = para.text.strip()
        if not line:
            continue
        all_lines.append(line)
        # Lấy định dạng từ run đầu tiên
        fmt = {}
        if para.runs:
            run = para.runs[0]
            fmt["bold"] = run.font.bold
            fmt["italic"] = run.font.italic
            fmt["font_size"] = run.font.size.pt if run.font.size else None
            if run.font.color and run.font.color.type:
                try:
                    fmt["color"] = str(run.font.color.rgb)
                except Exception:
                    pass
        paragraphs.append({
            "paragraph_no": p_idx,
            "text": line,
            "bullet_level": para.level,
            "formatting": fmt,
        })
    return "\n".join(all_lines), paragraphs

# ── Xác định vai trò block ───────────────────────────────────────
def detect_vai_tro(shape, slide_idx: int) -> str:
    name = (shape.name or "").lower()
    if "title" in name:
        return "tieu_de"
    if "subtitle" in name:
        return "phu_de"
    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
        return "hinh"
    try:
        if shape.table:  # sẽ raise nếu không phải table
            return "bang"
    except (ValueError, AttributeError):
        pass
    return "noi_dung"

# ── Trích từ khóa đơn giản từ text ──────────────────────────────
def extract_keywords(text: str) -> list[str]:
    # Lấy các từ có độ dài > 3, không trùng, ưu tiên chữ hoa đầu
    words = re.findall(r'[A-Za-zÀ-ỹà-ỹ]{4,}', text or "")
    seen = set()
    result = []
    for w in words:
        lw = w.lower()
        if lw not in seen:
            seen.add(lw)
            result.append(w)
        if len(result) >= 20:
            break
    return result

# ── Hàm import chính ─────────────────────────────────────────────
def import_pptx(filepath: Path, pg_conn, verbose=True) -> dict:
    """Import 1 file PPTX. Trả về dict thống kê."""
    stats = {"slides": 0, "blocks": 0, "texts": 0, "media": 0, "skipped": False}
    filepath = Path(filepath)

    # Kiểm tra trùng hash
    hash_val = file_hash(filepath)
    pc = pg_conn.cursor()
    pc.execute("SELECT ma_bai_giang FROM bai_giang WHERE ma_hash_tep = %s", (hash_val,))
    if pc.fetchone():
        if verbose:
            print(f"  ⏭  Bỏ qua (đã có): {filepath.name}")
        stats["skipped"] = True
        pc.close()
        return stats

    nguon_goc = detect_nguon_goc(filepath)
    prs = Presentation(str(filepath))
    total_slides = len(prs.slides)

    # Tạo bản ghi bai_giang
    pc.execute("""
        INSERT INTO bai_giang
            (tieu_de, nguon_goc, loai_tep, duong_dan_tep,
             ma_hash_tep, so_trang, ngay_them, trang_thai)
        VALUES (%s,%s,'pptx',%s,%s,%s,%s,'hoan_thanh')
        RETURNING ma_bai_giang
    """, (
        filepath.stem,
        nguon_goc,
        str(filepath),
        hash_val,
        total_slides,
        datetime.now(),
    ))
    ma_bg = pc.fetchone()[0]

    slide_titles = []
    slide_texts_all = []

    for s_idx, slide in enumerate(prs.slides):
        slide_title = ""
        slide_fulltext_parts = []
        layout_name = slide.slide_layout.name if slide.slide_layout else None

        # Tạo bản ghi trang
        pc.execute("""
            INSERT INTO trang
                (ma_bai_giang, so_thu_tu, tieu_de, loai_trang, ten_layout)
            VALUES (%s,%s,%s,'slide',%s)
            RETURNING ma_trang
        """, (ma_bg, s_idx + 1, None, layout_name))
        ma_trang = pc.fetchone()[0]

        block_order = 0
        for shape in slide.shapes:
            loai_khoi = "van_ban"
            vai_tro = detect_vai_tro(shape, s_idx)
            content_text = ""
            paragraphs = []

            if shape.has_text_frame:
                content_text, paragraphs = extract_shape_text(shape)
                if not content_text:
                    continue
                loai_khoi = "van_ban"
                if vai_tro == "tieu_de" and not slide_title:
                    slide_title = content_text
                slide_fulltext_parts.append(content_text)

            elif shape.shape_type == 13:  # PICTURE
                loai_khoi = "hinh_anh"

            else:
                # Kiểm tra table an toàn (pptx có thể raise ValueError)
                try:
                    tbl = shape.table
                    loai_khoi = "bang"
                    rows_data = []
                    for row in tbl.rows:
                        rows_data.append([cell.text.strip() for cell in row.cells])
                    pc.execute("""
                        INSERT INTO bang_du_lieu_trang
                            (ma_trang, so_hang, so_cot, hang_json)
                        VALUES (%s,%s,%s,%s)
                    """, (ma_trang, len(tbl.rows), len(tbl.columns), Json(rows_data)))
                    content_text = " | ".join(
                        [cell.text.strip() for row in tbl.rows for cell in row.cells if cell.text.strip()]
                    )
                except (ValueError, AttributeError):
                    pass  # Không phải table, bỏ qua

            # Tạo khoi_noi_dung
            pc.execute("""
                INSERT INTO khoi_noi_dung
                    (ma_trang, loai_khoi, vai_tro, noi_dung_van_ban, thu_tu_hien_thi)
                VALUES (%s,%s,%s,%s,%s)
                RETURNING ma_khoi
            """, (ma_trang, loai_khoi, vai_tro, content_text or None, block_order))
            ma_khoi = pc.fetchone()[0]
            block_order += 1
            stats["blocks"] += 1

            # doan_van_trang cho text shapes
            for p_info in paragraphs:
                pc.execute("""
                    INSERT INTO doan_van_trang
                        (ma_trang, ma_khoi, so_doan, cap_do_bullet,
                         vai_tro, noi_dung, noi_dung_chuan, dinh_dang_json)
                    VALUES (%s,%s,%s,%s,%s,%s,%s,%s)
                """, (
                    ma_trang, ma_khoi,
                    p_info["paragraph_no"],
                    p_info["bullet_level"],
                    vai_tro,
                    p_info["text"],
                    p_info["text"].lower(),
                    Json(p_info["formatting"]) if p_info["formatting"] else None,
                ))
                stats["texts"] += 1

        # Cập nhật tieu_de & noi_dung_van_ban cho trang
        fulltext = "\n".join(slide_fulltext_parts)
        pc.execute("""
            UPDATE trang SET tieu_de=%s, noi_dung_van_ban=%s WHERE ma_trang=%s
        """, (slide_title or None, fulltext or None, ma_trang))

        # tu_khoa_trang: upsert vào bang_tu_vung trước, rồi mới liên kết
        for kw in extract_keywords(fulltext):
            kw_clean = kw[:255]
            kw_lower = kw_clean.lower()
            # Upsert vào bang_tu_vung để lấy ma_tu_vung (UNIQUE constraint là tu_khoa)
            pc.execute("""
                INSERT INTO bang_tu_vung (tu_khoa, tu_khoa_chuan)
                VALUES (%s, %s)
                ON CONFLICT (tu_khoa) DO UPDATE SET tu_khoa_chuan = EXCLUDED.tu_khoa_chuan
                RETURNING ma_tu_vung
            """, (kw_clean, kw_lower))
            row = pc.fetchone()
            if row:
                ma_tv = row[0]
                pc.execute("""
                    INSERT INTO tu_khoa_trang (ma_trang, ma_tu_vung, nguon, trong_so)
                    VALUES (%s, %s, 'content', 1.0)
                    ON CONFLICT DO NOTHING
                """, (ma_trang, ma_tv))

        # chu_de_trang (cắt tối đa 255 ký tự)
        if slide_title:
            pc.execute("""
                INSERT INTO chu_de_trang (ma_trang, nhan_chu_de, loai_chu_de, thu_tu_hien_thi)
                VALUES (%s,%s,'title',%s)
            """, (ma_trang, slide_title[:255], s_idx + 1))

        slide_titles.append(slide_title)
        slide_texts_all.append(fulltext)
        stats["slides"] += 1

    # chu_de_bai_giang: lấy tiêu đề slide đầu tiên
    main_topic = slide_titles[0] if slide_titles else filepath.stem
    pc.execute("""
        INSERT INTO chu_de_bai_giang (ma_bai_giang, nhan_chu_de, loai_chu_de, thu_tu_hien_thi)
        VALUES (%s,%s,'title',1)
    """, (ma_bg, main_topic))

    # ho_so_bai_giang: tổng hợp nội dung
    all_text = "\n".join(slide_texts_all)
    pc.execute("""
        INSERT INTO ho_so_bai_giang (ma_bai_giang, chu_de_bai, noi_dung_bai)
        VALUES (%s,%s,%s)
        ON CONFLICT (ma_bai_giang) DO UPDATE SET chu_de_bai=EXCLUDED.chu_de_bai, noi_dung_bai=EXCLUDED.noi_dung_bai
    """, (ma_bg, main_topic, all_text[:2000] if all_text else None))

    # Cập nhật tieu_de bai_giang
    pc.execute("UPDATE bai_giang SET tieu_de=%s WHERE ma_bai_giang=%s",
               (main_topic or filepath.stem, ma_bg))

    pc.close()
    if verbose:
        print(f"  ✅ {filepath.name}: {stats['slides']} trang, {stats['blocks']} blocks, {stats['texts']} đoạn văn")
    return stats


# ── Batch import thư mục ─────────────────────────────────────────
def import_folder(folder: Path, recursive=True, verbose=True):
    """Import tất cả PPTX trong folder."""
    folder = Path(folder)
    pattern = "**/*.pptx" if recursive else "*.pptx"
    files = sorted(folder.glob(pattern))
    if not files:
        print(f"Không tìm thấy file PPTX trong: {folder}")
        return

    print(f"\n📂 Import thư mục: {folder}")
    print(f"   Tìm thấy {len(files)} file PPTX\n")

    pg = psycopg2.connect(**DB)
    pg.autocommit = False

    total = {"slides": 0, "blocks": 0, "texts": 0, "media": 0, "ok": 0, "skip": 0, "err": 0}

    for i, f in enumerate(files, 1):
        print(f"[{i}/{len(files)}] {f.name}")
        try:
            stats = import_pptx(f, pg, verbose=verbose)
            pg.commit()
            if stats["skipped"]:
                total["skip"] += 1
            else:
                total["ok"] += 1
                for k in ["slides", "blocks", "texts", "media"]:
                    total[k] += stats.get(k, 0)
        except Exception as e:
            pg.rollback()
            print(f"  ❌ Lỗi: {e}")
            import traceback; traceback.print_exc()
            total["err"] += 1

    pg.close()

    print(f"""
╔══════════════════════════════════════╗
║  KẾT QUẢ IMPORT PPTX                ║
╠══════════════════════════════════════╣
║  ✅ Thành công  : {total['ok']:>5} file           ║
║  ⏭  Bỏ qua     : {total['skip']:>5} file (đã có) ║
║  ❌ Lỗi         : {total['err']:>5} file           ║
║  📄 Trang/Slide : {total['slides']:>5}               ║
║  🧱 Khối ND     : {total['blocks']:>5}               ║
║  📝 Đoạn văn    : {total['texts']:>5}               ║
╚══════════════════════════════════════╝
""")


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Import PPTX vào kho_bai_giang")
    parser.add_argument("path", help="File .pptx hoặc thư mục chứa .pptx")
    parser.add_argument("--no-recurse", action="store_true", help="Không đệ quy thư mục con")
    args = parser.parse_args()

    p = Path(args.path)
    if p.is_file() and p.suffix.lower() == ".pptx":
        pg = psycopg2.connect(**DB)
        pg.autocommit = False
        import_pptx(p, pg)
        pg.commit()
        pg.close()
    elif p.is_dir():
        import_folder(p, recursive=not args.no_recurse)
    else:
        print("Lỗi: Đường dẫn không hợp lệ hoặc không phải .pptx")
        sys.exit(1)
